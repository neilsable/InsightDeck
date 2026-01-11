import os
os.environ["MPLCONFIGDIR"] = "/tmp/matplotlib"

from pathlib import Path
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from services.data_prep import load_and_validate_csv

def fit_picture_in_box(slide, image_path: str, x, y, box_w, box_h):
    img = Image.open(image_path)
    iw, ih = img.size
    img_ratio = iw / ih
    box_ratio = float(box_w) / float(box_h)

    if img_ratio >= box_ratio:
        target_w = box_w
        target_h = box_w / img_ratio
    else:
        target_h = box_h
        target_w = box_h * img_ratio

    x2 = x + (box_w - target_w) / 2
    y2 = y + (box_h - target_h) / 2

    slide.shapes.add_picture(image_path, x2, y2, width=target_w, height=target_h)

def make_chart(daily: pd.DataFrame, out_png: Path):
    df = daily.sort_values("day").copy()
    df["usage_7d"] = df["total_usage"].rolling(7, min_periods=1).mean()

    first = float(df["total_usage"].iloc[0])
    last = float(df["total_usage"].iloc[-1])
    change_pct = ((last - first) / first) * 100 if first else 0.0

    fig, ax = plt.subplots(figsize=(16, 5.2), dpi=200)

    ax.plot(df["day"], df["total_usage"], linewidth=1.1, alpha=0.35, label="Daily usage")
    ax.plot(df["day"], df["usage_7d"], linewidth=2.4, label="7-day rolling avg")

    ax.set_title("Cloud Platform Usage Trend", fontsize=15, pad=10, fontweight="bold")
    ax.set_xlabel("")
    ax.set_ylabel("Usage units", fontsize=11)

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(True, axis="y", linestyle="--", linewidth=0.8, alpha=0.18)
    ax.grid(False, axis="x")

    ax.xaxis.set_major_locator(mdates.WeekdayLocator(interval=1))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d %b"))
    plt.setp(ax.get_xticklabels(), fontsize=9)
    plt.setp(ax.get_yticklabels(), fontsize=9)

    x_last = df["day"].iloc[-1]
    y_last = df["total_usage"].iloc[-1]
    ax.scatter([x_last], [y_last], s=30, zorder=5)
    ax.annotate(
        f"Latest: {int(y_last):,}",
        xy=(x_last, y_last),
        xytext=(10, 10),
        textcoords="offset points",
        fontsize=9.5,
        fontweight="bold",
        bbox=dict(boxstyle="round,pad=0.30", alpha=0.10),
        arrowprops=dict(arrowstyle="->", alpha=0.35),
    )

    kpi_text = (
        f"Latest usage: {int(last):,} | Period change: {change_pct:+.1f}% | "
        f"Window: {df['day'].iloc[0].date()} → {df['day'].iloc[-1].date()}"
    )
    ax.text(0.01, 0.98, kpi_text, transform=ax.transAxes, va="top", fontsize=9.5, alpha=0.85)

    ax.legend(frameon=False, loc="upper right", fontsize=9)
    fig.tight_layout()
    fig.savefig(out_png)
    plt.close(fig)

def compute_metrics(df: pd.DataFrame) -> dict:
    daily = df.groupby("day", as_index=False).agg(
        total_usage=("usage_units", "sum"),
        total_cost=("cost_gbp", "sum"),
        total_incidents=("incidents", "sum"),
        avg_sla=("sla_pct", "mean"),
    ).sort_values("day")

    svc = df.groupby("service", as_index=False).agg(
        total_usage=("usage_units", "sum"),
        total_incidents=("incidents", "sum"),
        avg_sla=("sla_pct", "mean"),
    ).sort_values("total_usage", ascending=False)

    # period change
    usage_growth = ((daily["total_usage"].iloc[-1] - daily["total_usage"].iloc[0]) / daily["total_usage"].iloc[0]) * 100 if daily["total_usage"].iloc[0] else 0
    cost_growth = ((daily["total_cost"].iloc[-1] - daily["total_cost"].iloc[0]) / daily["total_cost"].iloc[0]) * 100 if daily["total_cost"].iloc[0] else 0

    kpis = {
        "usage_growth_pct": float(usage_growth),
        "cost_growth_pct": float(cost_growth),
        "sla_latest": float(daily["avg_sla"].iloc[-1]),
        "sla_overall": float(daily["avg_sla"].mean()),
        "incidents_total": int(daily["total_incidents"].sum()),
    }

    return {"daily": daily, "svc": svc, "kpis": kpis}

def build_narrative(kpis: dict, daily: pd.DataFrame, svc: pd.DataFrame) -> dict:
    # Week-over-week signal
    wow_line = "Momentum: insufficient history for week-over-week signal (needs ≥14 days)."
    if len(daily) >= 14:
        last7 = daily.tail(7)["total_usage"].mean()
        prev7 = daily.iloc[-14:-7]["total_usage"].mean()
        wow = ((last7 - prev7) / prev7) * 100 if prev7 else 0
        wow_line = f"Momentum: last 7-day avg {last7:,.0f} vs prior {prev7:,.0f} ({wow:+.1f}%)."

    # Service drivers
    drivers = []
    if not svc.empty:
        total_u = float(svc["total_usage"].sum()) or 1.0
        for _, r in svc.head(3).iterrows():
            share = (float(r["total_usage"]) / total_u) * 100
            drivers.append(
                f"{r['service']}: {share:.0f}% usage share, {int(r['total_incidents'])} incidents, SLA {float(r['avg_sla']):.3f}%"
            )

    insights = [
        f"Adoption: usage +{kpis['usage_growth_pct']:.1f}% across the period.",
        f"Spend: cost +{kpis['cost_growth_pct']:.1f}% (monitor cost-to-consumption).",
        f"Reliability: latest SLA {kpis['sla_latest']:.3f}% vs overall {kpis['sla_overall']:.3f}%.",
        wow_line,
        f"Operations: {kpis['incidents_total']} incidents logged across the period.",
    ]
    if drivers:
        insights.append("Top drivers: " + " | ".join(drivers[:2]))
        if len(drivers) > 2:
            insights.append("Additional driver: " + drivers[2])

    risks = []
    if kpis["sla_latest"] < 99.5:
        risks.append("SLA below threshold; customer impact risk if trend persists.")
    else:
        risks.append("SLA within tolerance; continue proactive monitoring and post-deploy checks.")

    if kpis["incidents_total"] > 120:
        risks.append("Incident volume elevated; risk of response fatigue and delivery drag.")
    else:
        risks.append("Incident volume manageable; keep RCA cadence and clear ownership for top drivers.")

    actions = [
        "Implement weekly cost-to-consumption governance and anomaly alerting.",
        "Run RCA on top incident drivers; add preventive checks in deployment gates.",
        "Introduce service-level SLOs per domain with clear owners and escalation paths.",
        "Publish a 5-minute exec snapshot weekly (WoW usage, SLA, incidents, cost).",
    ]

    return {"insights": insights, "risks": risks, "actions": actions}

def generate_ppt_from_csv(csv_path: Path, out_pptx: Path):
    df = load_and_validate_csv(csv_path)
    metrics = compute_metrics(df)
    daily = metrics["daily"]
    svc = metrics["svc"]
    kpis = metrics["kpis"]

    # Chart render
    out_png = Path("outputs") / f"chart_{out_pptx.stem}.png"
    out_png.parent.mkdir(exist_ok=True)
    make_chart(daily.rename(columns={"day": "day"}), out_png)

    narrative = build_narrative(kpis, daily, svc)

    # Theme
    BG = RGBColor(246, 248, 252)
    CARD = RGBColor(255, 255, 255)
    INK = RGBColor(15, 23, 42)
    MUTED = RGBColor(100, 116, 139)
    BORDER = RGBColor(226, 232, 240)

    ACCENT = RGBColor(37, 99, 235)
    ACCENT2 = RGBColor(16, 185, 129)
    VIOLET = RGBColor(99, 102, 241)
    CYAN = RGBColor(14, 165, 233)
    WARN = RGBColor(245, 158, 11)
    RISK = RGBColor(239, 68, 68)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    M = Inches(0.75)
    G = Inches(0.25)
    usable_w = SLIDE_W - 2 * M

    def card(sl, x, y, w, h, rounded=True):
        shp = sl.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = CARD
        shp.line.color.rgb = BORDER
        shp.line.width = Pt(1)
        return shp

    def textbox(sl, x, y, w, h, text="", size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, wrap=True):
        tb = sl.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return tb

    def badge(sl, x, y, w, h, text, fill, text_color=RGBColor(255, 255, 255)):
        b = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        b.fill.solid()
        b.fill.fore_color.rgb = fill
        b.line.color.rgb = fill
        tf = b.text_frame
        tf.clear()
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        return b

    # ----------------
    # Slide 1
    # ----------------
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()

    header_y = Inches(0.35); header_h = Inches(1.05)
    card(s1, M, header_y, usable_w, header_h)

    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, header_y, Inches(0.12), header_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    badge(s1, M + Inches(0.22), header_y + Inches(0.18), Inches(1.35), Inches(0.32), "OPS / INSIGHTS", ACCENT)
    textbox(s1, M + Inches(1.7), header_y + Inches(0.12), Inches(8.2), Inches(0.50), "Automated Executive KPI One-Pager", size=24, bold=True)
    textbox(s1, M + Inches(1.7), header_y + Inches(0.57), Inches(8.8), Inches(0.30), "Upload CSV → Auto-generate executive pack (synthetic-ready)", size=11, color=MUTED)

    badge_x = SLIDE_W - M - Inches(2.2)
    badge(s1, badge_x, header_y + Inches(0.18), Inches(2.2), Inches(0.32), "AUTO-GENERATED", ACCENT2)
    textbox(s1, badge_x, header_y + Inches(0.55), Inches(2.2), Inches(0.30), "Refresh: on upload  |  v1.0", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    tiles_y = header_y + header_h + Inches(0.28); tiles_h = Inches(1.28)
    tile_w = (usable_w - 3 * G) / 4

    sla_fill = ACCENT2 if kpis["sla_latest"] >= 99.7 else WARN if kpis["sla_latest"] >= 99.5 else RISK
    inc_fill = ACCENT2 if kpis["incidents_total"] <= 60 else WARN if kpis["incidents_total"] <= 90 else RISK

    def kpi_tile(ix, title, value, hint, strip_color, tag_text, tag_color):
        x = M + ix * (tile_w + G)
        card(s1, x, tiles_y, tile_w, tiles_h)

        strip = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, tiles_y, tile_w, Inches(0.08))
        strip.fill.solid(); strip.fill.fore_color.rgb = strip_color; strip.line.fill.background()

        textbox(s1, x + Inches(0.22), tiles_y + Inches(0.18), tile_w - Inches(0.44), Inches(0.22), title.upper(), size=9, bold=True, color=MUTED)
        textbox(s1, x + Inches(0.22), tiles_y + Inches(0.45), tile_w - Inches(0.44), Inches(0.48), value, size=22, bold=True)
        textbox(s1, x + Inches(0.22), tiles_y + Inches(0.98), tile_w - Inches(0.44), Inches(0.24), hint, size=10, color=MUTED)
        badge(s1, x + tile_w - Inches(1.05) - Inches(0.18), tiles_y + Inches(0.18), Inches(1.05), Inches(0.28), tag_text, tag_color)

    kpi_tile(0, "Usage Growth", f"+{kpis['usage_growth_pct']:.1f}%", "period change", ACCENT, "TREND", ACCENT)
    kpi_tile(1, "Cost Growth", f"+{kpis['cost_growth_pct']:.1f}%", "period change", VIOLET, "COST", VIOLET)
    kpi_tile(2, "SLA (Latest)", f"{kpis['sla_latest']:.3f}%", "target ≥ 99.7%", CYAN, "SLA", sla_fill)
    kpi_tile(3, "Incidents", f"{kpis['incidents_total']}", "total volume", RGBColor(168, 85, 247), "RISK", inc_fill)

    content_y = tiles_y + tiles_h + Inches(0.30)
    content_h = SLIDE_H - content_y - Inches(0.55)

    card(s1, M, content_y, usable_w, content_h)
    textbox(s1, M + Inches(0.30), content_y + Inches(0.20), usable_w - Inches(0.60), Inches(0.28), "Usage Trend", size=12, bold=True)
    textbox(s1, M + Inches(0.30), content_y + Inches(0.47), usable_w - Inches(0.60), Inches(0.22),
            "Chart auto-fitted to slide bounds (no overflow, aspect preserved).", size=10, color=MUTED)

    img_box_x = M + Inches(0.30)
    img_box_y = content_y + Inches(0.78)
    img_box_w = usable_w - Inches(0.60)
    img_box_h = content_h - Inches(1.35)

    fit_picture_in_box(s1, str(out_png), img_box_x, img_box_y, img_box_w, img_box_h)

    textbox(s1, M + Inches(0.30), content_y + content_h - Inches(0.40), usable_w - Inches(0.60), Inches(0.25),
            "Full narrative + drivers: see Appendix (Slide 2).", size=9, color=MUTED, align=PP_ALIGN.RIGHT, wrap=False)

    footer_y = SLIDE_H - Inches(0.42)
    textbox(s1, M, footer_y, usable_w, Inches(0.25),
            "Provenance: CSV → Python (metrics + chart + PPT) → Narrative engine (LLM swappable).",
            size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    # ----------------
    # Slide 2 (Dense Appendix)
    # ----------------
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg2.fill.solid(); bg2.fill.fore_color.rgb = BG; bg2.line.fill.background()

    header2_y = Inches(0.35); header2_h = Inches(0.90)
    card(s2, M, header2_y, usable_w, header2_h)
    textbox(s2, M + Inches(0.35), header2_y + Inches(0.18), Inches(11.5), Inches(0.5),
            "Appendix — Full Narrative & Drivers (Auto-generated)", size=22, bold=True)

    body_y = header2_y + header2_h + Inches(0.25)
    body_h = SLIDE_H - body_y - Inches(0.60)
    card(s2, M, body_y, usable_w, body_h)

    def section(title, title_fill, x, y, w, h, bullets):
        badge(s2, x, y, Inches(1.35), Inches(0.32), title, title_fill)
        tb = s2.shapes.add_textbox(x, y + Inches(0.42), w, h - Inches(0.48))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, it in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {it}"
            p.font.size = Pt(12)
            p.font.color.rgb = MUTED
            p.space_before = Pt(0)
            p.space_after = Pt(2)

    inner_x = M + Inches(0.40)
    inner_y = body_y + Inches(0.35)
    inner_w = usable_w - Inches(0.80)
    inner_h = body_h - Inches(0.55)

    col_gap = Inches(0.35)
    col_w = (inner_w - col_gap) / 2
    left_x = inner_x
    right_x = inner_x + col_w + col_gap

    section("INSIGHTS", ACCENT, left_x, inner_y, col_w, inner_h * 0.62, narrative["insights"][:10])
    section("RISKS", RISK, left_x, inner_y + inner_h * 0.65, col_w, inner_h * 0.33, narrative["risks"][:6])
    section("ACTIONS", ACCENT2, right_x, inner_y, col_w, inner_h * 0.62, narrative["actions"][:10])

    method = [
        "Input: CSV with day/service/usage/cost/incidents/SLA.",
        "Python computes KPIs + service drivers + WoW signals.",
        "Chart is auto-fitted into a bounding box (no overflow).",
        "Narrative deterministic today; can swap to LLM later.",
    ]
    section("METHOD", VIOLET, right_x, inner_y + inner_h * 0.65, col_w, inner_h * 0.33, method)

    out_pptx.parent.mkdir(exist_ok=True)
    prs.save(out_pptx)


