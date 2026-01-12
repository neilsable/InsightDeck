from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Tuple

import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


REQUIRED_COLS = {"day", "service", "usage_units", "cost_gbp", "incidents", "sla_pct"}

# ===== Theme =====
@dataclass(frozen=True)
class Theme:
    bg: RGBColor = RGBColor(18, 32, 56)          # deep navy
    panel: RGBColor = RGBColor(28, 48, 78)       # slightly lighter panel
    card: RGBColor = RGBColor(245, 247, 250)     # light cards
    card_line: RGBColor = RGBColor(220, 225, 232)
    text: RGBColor = RGBColor(245, 249, 255)     # white-ish
    subtext: RGBColor = RGBColor(185, 205, 230)  # muted
    ink: RGBColor = RGBColor(25, 28, 35)         # dark text on cards
    accent: RGBColor = RGBColor(64, 132, 255)    # blue
    accent2: RGBColor = RGBColor(0, 205, 160)    # green
    warn: RGBColor = RGBColor(255, 186, 0)       # amber
    danger: RGBColor = RGBColor(255, 90, 90)     # red

THEME = Theme()

# ===== Layout Grid (16:9) =====
# PowerPoint default widescreen is 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_X = Inches(0.8)
MARGIN_Y = Inches(0.6)
GAP = Inches(0.25)

FONT_TITLE = "Aptos Display"   # fallback ok if unavailable
FONT_BODY = "Aptos"            # fallback ok if unavailable

# ===== Helpers: formatting & safe text =====
def _format_gbp(x: float) -> str:
    if x >= 1_000_000:
        return f"£{x/1_000_000:.2f}M"
    if x >= 1_000:
        return f"£{x/1_000:.1f}K"
    return f"£{x:.0f}"

def _safe_lines(lines: List[str], max_lines: int, max_chars: int) -> List[str]:
    """Hard safety: prevent overflow by capping lines and truncating."""
    out = []
    for s in lines[:max_lines]:
        s = s.strip()
        if len(s) > max_chars:
            s = s[: max_chars - 1].rstrip() + "…"
        out.append(s)
    return out

def _fit_text(tf, lines: List[str], start_size: int, min_size: int, max_lines: int, max_chars: int):
    """Soft safety: shrink font if needed; then truncate."""
    tf.clear()
    tf.word_wrap = True

    safe = _safe_lines(lines, max_lines=max_lines, max_chars=max_chars)

    # start with a size and add paragraphs
    size = start_size
    while True:
        tf.clear()
        first = True
        for ln in safe:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = ln
            p.font.size = Pt(size)
            p.font.name = FONT_BODY

        # Heuristic: if size is too big for many lines, shrink
        # (python-pptx has no perfect text measurement; this pragmatic rule works well.)
        if len(safe) <= 3 or size <= min_size:
            break
        if len(" ".join(safe)) > (max_chars * max_lines * 0.75) and size > min_size:
            size -= 1
        else:
            break

    # Final truncate if still too dense
    if size == min_size and len(safe) == max_lines:
        # keep as-is; already truncated by chars
        pass

def _add_bg(slide, theme: Theme = THEME):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = theme.bg

def _add_header(slide, title: str, subtitle: str):
    # Title
    tb = slide.shapes.add_textbox(MARGIN_X, MARGIN_Y, SLIDE_W - 2*MARGIN_X, Inches(0.7))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = THEME.text

    # Subtitle
    sb = slide.shapes.add_textbox(MARGIN_X, MARGIN_Y + Inches(0.65), SLIDE_W - 2*MARGIN_X, Inches(0.4))
    stf = sb.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.name = FONT_BODY
    sp.font.size = Pt(16)
    sp.font.color.rgb = THEME.subtext

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_X, MARGIN_Y + Inches(1.15),
        Inches(2.2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME.accent
    line.line.fill.background()

def _add_panel(slide, x, y, w, h):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = THEME.panel
    panel.line.color.rgb = RGBColor(30, 45, 70)
    panel.line.width = Pt(1)
    return panel

def _add_card(slide, x, y, w, h, label: str, value: str, accent: RGBColor):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = THEME.card
    card.line.color.rgb = THEME.card_line
    card.line.width = Pt(1)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Label
    lb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.18), w - Inches(0.4), Inches(0.3))
    ltf = lb.text_frame
    ltf.clear()
    p1 = ltf.paragraphs[0]
    p1.text = label.upper()
    p1.font.name = FONT_BODY
    p1.font.size = Pt(11)
    p1.font.color.rgb = RGBColor(95, 105, 120)

    # Value
    vb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.50), w - Inches(0.4), Inches(0.6))
    vtf = vb.text_frame
    vtf.clear()
    p2 = vtf.paragraphs[0]
    p2.text = value
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = THEME.ink

def _chart_cost_trend(df: pd.DataFrame, out_png: Path):
    daily = df.groupby("day", as_index=False)["cost_gbp"].sum()

    fig = plt.figure(figsize=(10, 3.2))
    ax = plt.gca()

    # Line
    ax.plot(daily["day"], daily["cost_gbp"], linewidth=2.6, marker="o", markersize=4)

    # Make ALL chart text visible on dark backgrounds
    title_color = "#F5F9FF"
    tick_color = "#D7E3F7"
    grid_color = "#9DB2D6"

    ax.set_title("Cloud Cost Trend (GBP)", fontsize=13, color=title_color, pad=10)
    ax.set_xlabel("")
    ax.set_ylabel("")

    ax.tick_params(axis="x", colors=tick_color, labelsize=9)
    ax.tick_params(axis="y", colors=tick_color, labelsize=9)

    ax.grid(True, axis="y", alpha=0.22, color=grid_color)

    # Spines
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_alpha(0.25)
    ax.spines["bottom"].set_alpha(0.25)
    ax.spines["left"].set_color(grid_color)
    ax.spines["bottom"].set_color(grid_color)

    # Transparent background (so panel shows through)
    fig.patch.set_alpha(0.0)
    ax.set_facecolor("none")

    plt.tight_layout()
    plt.savefig(out_png, dpi=200, transparent=True)
    plt.close()


@dataclass
class KPIs:
    total_cost: float
    total_usage: float
    avg_sla: float
    total_incidents: int

def _validate_load(csv_path: Path) -> pd.DataFrame:
    df = pd.read_csv(csv_path)
    missing = REQUIRED_COLS - set(df.columns)
    if missing:
        raise ValueError(f"Missing required columns: {sorted(missing)}")

    df["day"] = pd.to_datetime(df["day"], errors="coerce")
    if df["day"].isna().any():
        raise ValueError("Invalid date format in 'day'. Use YYYY-MM-DD.")

    for c in ["usage_units", "cost_gbp", "incidents", "sla_pct"]:
        df[c] = pd.to_numeric(df[c], errors="coerce")
    if df[["usage_units", "cost_gbp", "incidents", "sla_pct"]].isna().any().any():
        raise ValueError("Numeric columns contain invalid values.")

    return df.sort_values("day")

def _kpis(df: pd.DataFrame) -> KPIs:
    return KPIs(
        total_cost=float(df["cost_gbp"].sum()),
        total_usage=float(df["usage_units"].sum()),
        avg_sla=float(df["sla_pct"].mean()),
        total_incidents=int(df["incidents"].sum()),
    )

def generate_ppt_from_csv(csv_path: Path, out_pptx: Path) -> None:
    csv_path = Path(csv_path)
    out_pptx = Path(out_pptx)

    df = _validate_load(csv_path)
    k = _kpis(df)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # === Slide 1 ===
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(s1)
    _add_header(s1, "InsightDeck Executive Summary", "KPI one-pager + auto-fitted trend chart")

    # KPI row area
    cards_y = MARGIN_Y + Inches(1.45)
    card_h = Inches(1.05)
    card_w = (SLIDE_W - 2*MARGIN_X - 3*GAP) / 4

    _add_card(s1, MARGIN_X + (card_w+GAP)*0, cards_y, card_w, card_h, "Total Cost", _format_gbp(k.total_cost), THEME.accent)
    _add_card(s1, MARGIN_X + (card_w+GAP)*1, cards_y, card_w, card_h, "Total Usage", f"{k.total_usage:,.0f}", THEME.accent2)
    _add_card(s1, MARGIN_X + (card_w+GAP)*2, cards_y, card_w, card_h, "Avg SLA", f"{k.avg_sla:.2f}%", THEME.warn)
    _add_card(s1, MARGIN_X + (card_w+GAP)*3, cards_y, card_w, card_h, "Incidents", f"{k.total_incidents:,}", THEME.danger)

    # Chart panel
    chart_y = cards_y + card_h + Inches(0.35)
    chart_h = SLIDE_H - chart_y - MARGIN_Y
    panel = _add_panel(s1, MARGIN_X, chart_y, SLIDE_W - 2*MARGIN_X, chart_h)

    chart_path = Path("/tmp") / f"insightdeck_cost_trend_{out_pptx.stem}.png"
    _chart_cost_trend(df, chart_path)

    # Put chart image inside panel with padding
    pad = Inches(0.3)
    s1.shapes.add_picture(
        str(chart_path),
        MARGIN_X + pad,
        chart_y + pad,
        width=(SLIDE_W - 2*MARGIN_X - 2*pad),
        height=(chart_h - 2*pad)
    )

    # === Slide 2 ===
    s2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(s2)
    _add_header(s2, "Insights, Risks & Actions", "Narrative appendix (auto-generated)")

    # Build narrative safely
    daily = df.groupby("day", as_index=False).agg(
        cost_gbp=("cost_gbp", "sum"),
        incidents=("incidents", "sum"),
        sla_pct=("sla_pct", "mean"),
        usage_units=("usage_units", "sum"),
    )
    cost_change = float(daily["cost_gbp"].iloc[-1] - daily["cost_gbp"].iloc[0]) if len(daily) > 1 else 0.0
    trend = "up" if cost_change > 0 else "down" if cost_change < 0 else "flat"

    insights = [
        f"Cost trend is {trend} over the period ({_format_gbp(abs(cost_change))} net change).",
        f"Average SLA is {k.avg_sla:.2f}%, with {k.total_incidents} total incidents.",
        f"Peak daily cost: {_format_gbp(float(daily['cost_gbp'].max()))}.",
    ]
    risks = [
        "Rising cost with flat usage can indicate inefficiency or pricing drift.",
        "Incident spikes may correlate with SLA degradation and operational risk.",
    ]
    actions = [
        "Review top-cost services; apply rightsizing and usage caps where appropriate.",
        "Investigate incident spikes; add alerts, SLOs, and runbooks.",
        "Use InsightDeck weekly to standardize stakeholder reporting.",
    ]
    methods = [
        "Input validated for required columns and numeric types.",
        "Daily aggregation used for trend chart and KPI rollups.",
        "Deck generated using python-pptx; charts rendered via matplotlib (Agg).",
    ]

    # Two-column grid
    grid_top = MARGIN_Y + Inches(1.55)
    col_w = (SLIDE_W - 2*MARGIN_X - GAP) / 2
    row_h = Inches(2.35)

    def add_block(title: str, lines: List[str], x, y, accent: RGBColor):
        # Panel
        _add_panel(s2, x, y, col_w, row_h)
        # Title
        tb = s2.shapes.add_textbox(x + Inches(0.35), y + Inches(0.25), col_w - Inches(0.7), Inches(0.35))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = THEME.text

        # Accent
        bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.35), y + Inches(0.62), Inches(1.4), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Body (fit-safe)
        body = s2.shapes.add_textbox(x + Inches(0.35), y + Inches(0.8), col_w - Inches(0.7), row_h - Inches(1.0))
        btf = body.text_frame
        btf.vertical_anchor = MSO_ANCHOR.TOP
        btf.word_wrap = True
        lines2 = _safe_lines(lines, max_lines=5, max_chars=92)
        # render bullets
        btf.clear()
        for i, ln in enumerate(lines2):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = ln
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.color.rgb = THEME.text

    add_block("Key Insights", insights, MARGIN_X, grid_top, THEME.accent)
    add_block("Key Risks", risks, MARGIN_X + col_w + GAP, grid_top, THEME.warn)
    add_block("Recommended Actions", actions, MARGIN_X, grid_top + row_h + GAP, THEME.accent2)
    add_block("Method Notes", methods, MARGIN_X + col_w + GAP, grid_top + row_h + GAP, RGBColor(130, 150, 180))

    prs.save(out_pptx)
